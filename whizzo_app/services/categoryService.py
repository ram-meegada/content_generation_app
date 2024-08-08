from langdetect import detect
from datetime import datetime, timedelta
import ast
from ssl import SSL_ERROR_EOF
from typing import final
from django.http import JsonResponse
from whizzo_app.models.assignmentModel import AssignmentModel
from whizzo_app.models.testingModel import TestingModel
from whizzo_app.models.achievementModel import AchievementModel
from whizzo_app.models.abilityModel import AbilityModel
from whizzo_app.models.categoryModel import CategoryModel
from whizzo_project import settings
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_core.messages import HumanMessage
from dotenv import load_dotenv
import textwrap
import urllib.request as urlopener
from PyPDF2 import PdfReader
from io import BytesIO
from whizzo_app.models import FaqModel, CmsModel, UserModel, FileSumarizationModel, NoteModel, TestingModel, PresentationModel, NoteTakingModel
from whizzo_app.utils import messages
from whizzo_app.serializers import categorySerializer, adminSerializer
from deep_translator import GoogleTranslator
import aspose.words as aw
from pdf2docx import Converter
import random
from whizzo_app.utils.saveImage import save_file_conversion
from whizzo_app.serializers.uploadMediaSerializer import CreateUpdateUploadMediaSerializer
import os
import tabula
import re
import json
import xlsxwriter
import pandas as pd
from reportlab.lib.pagesizes import letter, A3
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Image
from reportlab.lib import colors
from pdf2image import convert_from_path
from django.core.files.storage import FileSystemStorage
import speech_recognition as sr
from pydub import AudioSegment
from whizzo_app.utils.saveImage import save_image
from whizzo_app.utils import sendMail
from whizzo_app.services.uploadMediaService import UploadMediaService
from whizzo_app.utils.customPagination import CustomPagination
from docx import Document
from docx2pdf import convert
import fitz
from PIL import Image
import tempfile
from pptx.util import Inches
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from rest_framework import status
import pdfkit
import pypandoc
from whizzo_app.utils.saveImage import saveFile
from whizzo_app.services.uploadMediaService import UploadMediaService
from whizzo_app.utils.sendMail import send_pdf_file_to_mail
from threading import Thread
import os
from io import BytesIO
from django.core.files import File
import pytesseract
from whizzo_app.models.fileConversionModel import FileConversationModel
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from PIL import Image
from PIL import Image, ImageDraw
from googletrans import Translator
from django.core.files.uploadedfile import UploadedFile
# from spire.presentation.common import *
# from spire.presentation import *
from bs4 import BeautifulSoup
from whizzo_app.models.articleModel import ArticleModel
from whizzo_app.utils.Modules.testingModule import chatGPT_pdf_processing, extract_data_from_url, chatGPT_image_processing
from whizzo_app.utils.Modules.articlesModule import generate_article_util
from whizzo_app.utils.Modules.assignmentSolutionsModule import assigment_chatGPT_pdf_processing, assignment_extract_text
from whizzo_app.utils.Modules.presentationModule import generate_presentation_util

pytesseract.pytesseract.tesseract_cmd = 'C:\\Program Files\\Tesseract-OCR\\tesseract.exe'

upload_media_obj = UploadMediaService()

load_dotenv()
google_api_key = settings.GOOGLE_API_KEY


def generate_file_name(name):
    import string
    char = string.ascii_letters + string.digits + "."
    for i in name:
        if i not in char:
            name = name.replace(i, '')
    try:
        for j in range(len(name)-1, -1, -1):
            if name[j] == ".":
                extension = name[j+1:]
                file_name = name[:j]
                break
        result = [file_name, extension]
        return result
    except:
        return [name, None]


def to_markdown(text):
    text = text.replace('*', '').replace('#', '').replace("-", "")
    intent_text = (textwrap.indent(text, '', predicate=lambda _: True))
    return intent_text

# def to_markdown(text):
#     # Remove Markdown headers (e.g., # Header)
#     text = re.sub(r'^\s*#+\s+', '', text, flags=re.MULTILINE)
#     # Remove emphasis (e.g., *italic* or **bold**)
#     text = re.sub(r'(\*|_){1,2}(.*?)\1{1,2}', r'\2', text)
#     # Remove inline code (e.g., `code`)
#     text = re.sub(r'`(.*?)`', r'\1', text)
#     # Remove strikethrough (e.g., ~~text~~)
#     text = re.sub(r'~~(.*?)~~', r'\1', text)
#     # Remove links (e.g., [text](url))
#     text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
#     # Remove images (e.g., ![alt text](url))
#     text = re.sub(r'!\[(.*?)\]\(.*?\)', r'\1', text)
#     # Remove remaining special characters used in Markdown
#     text = re.sub(r'[*_~`]', '', text)

#     # Normalize indentation
#     text = textwrap.dedent(text)

#     # Split text into lines and clean up
#     lines = [line.strip() for line in text.splitlines() if line.strip()]

#     # Join lines to form a single JSON string
#     json_string = " ".join(lines)

#     # Attempt to parse as JSON
#     try:
#         json_list = json.loads(json_string)
#     except json.JSONDecodeError as e:
#         raise ValueError(f"Failed to decode JSON: {e}")

#     return json_list

# def to_markdown(text):
#     # Remove Markdown headers (e.g., # Header)
#     text = re.sub(r'^\s*#+\s+', '', text, flags=re.MULTILINE)
#     # Remove emphasis (e.g., *italic* or **bold**)
#     text = re.sub(r'(\*|_){1,2}(.*?)\1{1,2}', r'\2', text)
#     # Remove inline code (e.g., `code`)
#     text = re.sub(r'`(.*?)`', r'\1', text)
#     # Remove strikethrough (e.g., ~~text~~)
#     text = re.sub(r'~~(.*?)~~', r'\1', text)
#     # Remove links (e.g., [text](url))
#     text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
#     # Remove images (e.g., ![alt text](url))
#     text = re.sub(r'!\[(.*?)\]\(.*?\)', r'\1', text)
#     # Remove remaining special characters used in Markdown
#     text = re.sub(r'[*_~`]', '', text)

#     # Normalize indentation
#     text = textwrap.dedent(text)

#     # Split text into lines and clean up
#     lines = [line.strip() for line in text.splitlines() if line.strip()]

#     # Convert lines to JSON list
#     try:
#         json_list = json.loads("[" + ",".join(lines) + "]")
#     except json.JSONDecodeError:
#         json_list = lines  # Fallback: return lines as plain list if JSON decoding fails

#     return json_list


def assignment_to_markdown(text):
    # Remove Markdown headers (e.g., # Header)
    text = re.sub(r'^\s*#+\s+', '', text, flags=re.MULTILINE)
    # Remove emphasis (e.g., *italic* or **bold**)
    text = re.sub(r'(\*|_){1,2}(.*?)\1{1,2}', r'\2', text)
    # Remove inline code (e.g., `code`)
    text = re.sub(r'`(.*?)`', r'\1', text)
    # Remove strikethrough (e.g., ~~text~~)
    text = re.sub(r'~~(.*?)~~', r'\1', text)
    # Remove links (e.g., [text](url))
    text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
    # Remove images (e.g., ![alt text](url))
    text = re.sub(r'!\[(.*?)\]\(.*?\)', r'\1', text)
    # Remove remaining special characters used in Markdown
    text = re.sub(r'[*_~`]', '', text)

    # Normalize indentation
    text = textwrap.dedent(text)

    # Split text into lines and clean up
    lines = [line.strip() for line in text.splitlines() if line.strip()]

    # Remove commas from each line
    lines = [line.replace(',', '') for line in lines]

    # Convert lines to JSON list
    try:
        json_list = json.loads("[" + ",".join(lines) + "]")
    except json.JSONDecodeError:
        json_list = lines  # Fallback: return lines as plain list if JSON decoding fails

    return json_list


def image_processing(image_link, query):
    '''processing the image and generate the mcq with options and answer 
       image_link is the s3 bucket link of image and query is string 
       which we use to generate the mcq of flashcards'''
    llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash")
    # example
    message = HumanMessage(
        content=[
            {
                "type": "text",
                "text": query,
            },
            {"type": "image_url", "image_url": str(image_link)},
        ]
    )
    response = llm.invoke([message])
    result_qu = to_markdown(response.content)
    return result_qu


def pdf_processing(pdf_file, query):
    '''processing the pdf and generate the mcq and flash cards based on the user requirement 
       pdf_link is the s3 bucket link where we store the pdf file 
       query is based on the user selection if he select mcq then query based on mcq else based on flash card'''

    llm = ChatGoogleGenerativeAI(model="gemini-pro")

    def extract_text(pdf_file):
        '''extract all the text from the pdf '''
        pdf_text = ''
        # Wrap the file content in BytesIO to treat it as a bytes-like object
        url_obj = urlopener.urlopen(pdf_file)

        # with open(pdf_file, 'r') as f:
        pdf_stream = BytesIO(url_obj.read())
        pdf_reader = PdfReader(pdf_stream)
        for page in pdf_reader.pages:
            pdf_text += page.extract_text()
        print(pdf_text, '------pdf oodfsd')
        return pdf_text

    if pdf_file:
        try:
            text_data = extract_text(pdf_file)
            if text_data.strip() == "":
                return {"message": "Empty file provided."}
            message = HumanMessage(
                content=[
                    {"type": "text",
                     "text": query},
                    {"type": "text", "text": text_data}
                ]
            )

            # Process the message
            response = llm.invoke([message])
            result = to_markdown(response.content)
            return result
        except Exception as e:
            return (str(e))
    else:
        return ("No PDF file provided")


class CategoryService:
    def generate_testing_category_result(self, request):
        try:
            file_links = []
            files = dict(request.data)["file"]
            for i in files:
                file_links.append(save_image(i)[0])
            print(file_links, '--------')    
            sub_category = int(request.data["sub_category"])
            api_type = int(request.data["type"])
            final_response = []
            if sub_category == 1:
                if api_type == 1:
                    number_of_questions = int(
                        settings.NUMBER_OF_QUESTIONS)//len(file_links)
                    for file in file_links:
                        query = f"Generate {number_of_questions} mcqs with options and answers for this image and make in python json list format. Keys should be 'question_no', 'question', 'answer_option', 'correct_answer'.)"
                        result = chatGPT_image_processing(file, query)
                        final_response += result
                elif api_type == 2:
                    text_data = extract_data_from_url(file_links[0])
                    if "ar" in detect(text_data[0]):
                        input_language = "arabic"
                    else:
                        input_language = "english"
                    number_of_questions = int(
                        settings.NUMBER_OF_QUESTIONS)//len(text_data)
                    query = f"Generate minimum of {number_of_questions} mcqs with options and answers for the input in {input_language} language. Format should be in python json list of objects(key-value pair). Key names of objects should be strictly 'question_no', 'question', 'answer_option' and 'correct_answer'."
                    for i in text_data:
                        result = chatGPT_pdf_processing(i, query)
                        for j in result:
                            final_response.append(j)
                    print(final_response, '-----final response-----')
                    for index, body in enumerate(final_response):
                        if "answer_options" in body:
                            body["answer_option"] = body.pop("answer_options")
                        if isinstance(body["answer_option"], dict):
                            if "correct_answer" not in body and "correct_answer" in body["answer_option"]:
                                correct_answer = body["answer_option"].pop(
                                    "correct_answer")
                                body["correct_answer"] = correct_answer
                            if "correct_answer" in body:
                                for key, val in body["answer_option"].items():
                                    if key == body["correct_answer"]:
                                        body["correct_answer"] = val
                            body["answer_option"] = list(
                                body["answer_option"].values())
                        body["question_no"] = index + 1
            elif sub_category == 2:
                if api_type == 1:
                    number_of_questions = int(
                        settings.NUMBER_OF_QUESTIONS)//len(file_links)
                    for file in file_links:
                        query = f"Generate {number_of_questions} flashcards for this input. Format should be in python json list. Keys should be 'question', 'answer'. 'question' and 'answer' should be different. if the question has multiple answers give them in list."
                        result = chatGPT_image_processing(file, query)
                        final_response += result
                elif api_type == 2:
                    text_data = extract_data_from_url(file_links[0])
                    number_of_questions = int(
                        settings.NUMBER_OF_QUESTIONS)//len(text_data)
                    query = f"First find the language of input and Generate {number_of_questions} flashcards for this input in same language. Format should be in python json list. Keys should be 'question', 'answer'. if the question has multiple answers give them in list."
                    for i in text_data:
                        result = chatGPT_pdf_processing(i, query)
                        for i in result:
                            final_response.append(i)
                for i in final_response:
                    if isinstance(i["answer"], str):
                        i["answer"] = [i["answer"]]
            final_response = [i for i in final_response if isinstance(i, dict)]
            save_data = TestingModel.objects.create(user_id=request.user.id,
                                                    sub_category=request.data["sub_category"],
                                                    result=final_response,
                                                    remaining_answers=len(
                                                        final_response),
                                                    sub_category_type=sub_category
                                                    )
            print(final_response, '-----final_response-----final_response----')                                        
            if not final_response:
                return {"data": None, "message": "Please upload again", "status": 400}
            return {"data": final_response, "record_id": save_data.id, "message": "Result generated successfully", "status": 200}
        except Exception as error:
            print(error, '-----error--------')
            return {"data": str(error), "message": "Something went wrong", "status": 400}

    def generate_testing_category_result_pdf(self, request):
        file_link = request.data["file_link"]
        sub_category = request.data["sub_category"]
        final_response = []
        if sub_category == 1:
            query = f"Generate {settings.NUMBER_OF_QUESTIONS} mcqs with options and answers for this input. Format should be in python json list."
            result = pdf_processing(file_link, query)
            json_result = self.jsonify_response(result)
            final_response.append(json_result)
        elif sub_category == 2:
            query = f"Generate {settings.NUMBER_OF_QUESTIONS} flashcards for this input. Format should be in python json list."
            result = pdf_processing(file_link, query)
            json_result = self.jsonify_response(result)
            final_response.append(json_result)
        create_category = CategoryModel.objects.create(user_id=request.user.id,
                                                       category=1,  # static because this api is for testing category
                                                       sub_category=request.data["sub_category"],
                                                       result=final_response
                                                       )
        return {"data": final_response, "message": "Result generated successfully", "status": 200}

    def jsonify_response(self, result):
        for i, j in enumerate(result):
            if j == "[":
                break
        for end in range(len(result)-1, -1, -1):
            if result[end] == "]":
                break
        final_result = json.loads(result[i:end+1])
        for i in final_result:
            i["user_answer"] = ""
        return final_result

    def submit_test_and_update_result(self, request, id):
        get_test_object = TestingModel.objects.get(id=id)
        user_response = request.data.get("user_response")
        print(user_response, '------user response-------')
        correct_answers, wrong_answers, remaining_answers = 0, 0, 0
        if get_test_object.sub_category == 1 or (get_test_object.sub_category in [3, 4] and get_test_object.sub_category_type in [1]):
            for i in user_response:
                if i["correct_answer"] == i["user_answer"]:
                    correct_answers += 1
                elif i["user_answer"] == "":
                    remaining_answers += 1
                elif i["correct_answer"] != i["user_answer"]:
                    wrong_answers += 1
        elif get_test_object.sub_category in [2] or (get_test_object.sub_category in [3, 4] and get_test_object.sub_category_type in [2]):
            for i in user_response:
                if i["user_answer"] == "YES":
                    correct_answers += 1
                elif i["user_answer"] == "":
                    remaining_answers += 1
                elif i["user_answer"] == "NO":
                    wrong_answers += 1
        get_test_object.correct_answers = correct_answers
        get_test_object.wrong_answers = wrong_answers
        get_test_object.remaining_answers = remaining_answers
        # get_test_object.result = user_response
        get_test_object.save()
        total_questions = len(user_response)
        correct_answers_percentage = round(
            (correct_answers/total_questions)*100, 2)
        wrong_answers_percentage = round(
            (wrong_answers/total_questions)*100, 2)
        remaining_answers_percentage = round(
            (remaining_answers/total_questions)*100, 2)

        data = {"correct_answers": correct_answers, "wrong_answers": wrong_answers, "remaining_answers": remaining_answers,
                "correct_answers_percentage": correct_answers_percentage, "wrong_answers_percentage": wrong_answers_percentage,
                "remaining_answers_percentage": remaining_answers_percentage}
        return {"data": data, "message": messages.TEST_SUBMITTED, "status": 200}

    def previous_tests_listing(self, request):
        if "sub_category" not in request.data:
            previous_tests = TestingModel.objects.filter(
                user_id=request.user.id).order_by("-updated_at")
        elif "sub_category" in request.data:
            previous_tests = TestingModel.objects.filter(
                user_id=request.user.id, sub_category__in=request.data["sub_category"]).order_by("-updated_at")
        pagination_obj = CustomPagination()
        search_keys = []
        result = pagination_obj.custom_pagination(
            request, search_keys, categorySerializer.GetPreviousTestSerializer, previous_tests)
        return {"data": result, "message": messages.TESTING_CATEGORY_PAST_TESTS, "status": 200}


# filesumarization


    def generate_file_summary(self, request):
        llm = ChatGoogleGenerativeAI(model="gemini-pro")
        if int(request.data["type"]) == 1:
            file_link = request.FILES.get("file_link")
            try:
                text_data = self.extract_text(file_link)
                print(text_data, '------ text data --------')
                message = HumanMessage(
                    content=[
                        {"type": "text",
                         "text": f"generate a summary of the input I provide you and the length of the summary should be strictly atleast 2000 words and give me only text no * and extra symbols"},
                        {"type": "text", "text": text_data}
                    ]
                )
                response = llm.invoke([message])
                result = to_markdown(response.content)
                # message=[
                #             {"role": "system", "content": "You are a helpful assistant designed to output JSON."},
                #             {"role": "user", "content": [
                #                     {"type": "text", "text": "Generate a summary of the input I provide you and the length of the summary should be strictly atleast 2000 words. Please maintain line breaks and intendations."},
                #                     {"type": "text", "text": f"Input data: {text_data}"}
                #                 ]}
                #             ]
                # chatbot = openai.ChatCompletion.create(
                #     model="gpt-3.5-turbo", messages=message,response_format={ "type": "json_object" },temperature = 0.0,
                # )
                # reply = chatbot.choices[0].message.content
                # try:
                #     final_data = json.loads(reply)
                # except:
                #     final_data = ast.literal_eval(reply)    
                # print(final_data, '----- text data ----------- text data ------')
                # result = list(final_data.values())[0]
                save_file_summary_record = FileSumarizationModel.objects.create(
                    user_id=request.user.id,
                    sub_category=5,
                    result=result
                )
                # Store the result in the session or a temporary variable
                return {"data": result, "record_id": save_file_summary_record.id, "message": messages.SUMMARY_GENERATED, "status": 200}
            except Exception as e:
                return {"error": str(e), "message": messages.PLEASE_UPLOAD_AGAIN, "status": 400}
        elif int(request.data["type"]) == 2:
            images = dict(request.data)["file_link"]
            query = "You are summary provider. Generate a summary of text I provide you and the length of the summary should be strictly atleast 2000 words and give me only text no * and extra symbols"
            try:
                text_data = ""
                gemini_result = []
                for file_image in images:
                    img = save_image(file_image)
                    result = self.image_processing_assignment_solution(
                        img[0], query)
                    # text_data = self.extract_text_from_image(file_image)
                    # message = HumanMessage(
                    #     content=[
                    #         {"type": "text",
                    #             "text": "You are a teacher. Generate questions and answers based on the data I provide to you. Format should be proper Python Javascript object notation list of dictionaries where every dictionary contains keys as 'question', 'answer' and 'options'(if available)."},
                    #         {"type": "text", "text":text_data}
                    #     ]
                    # )
                    # # result = self.gemini_solution(file_link)
                    # response = llm.invoke([message])
                    # result = to_markdown(response.content)
                    temp = self.format_final_response(result)
                    gemini_result += temp
                save_file_summary_record = FileSumarizationModel.objects.create(
                    user_id=request.user.id,
                    sub_category=5,
                    result=result
                )
                try:
                    self.translate_text_for_file_summarization(
                        result, save_file_summary_record)
                except:
                    pass
                return {"data": result, "record_id": save_file_summary_record.id, "message": messages.SUMMARY_GENERATED, "status": 200}
            except Exception as err:
                return {"error": str(err), "message": messages.PLEASE_UPLOAD_AGAIN, "status": 400}

    def translate_text_for_file_summarization(self, text, save_file_summary_record):
        text_list = text.split(" ")
        translator = Translator()
        translations = []
        for text in text_list:
            translation = translator.translate(text, src="en", dest="ar")
            translations.append(translation.text)
        save_file_summary_record.arabic_result = " ".join(translations)
        save_file_summary_record.save()
        return None

    def extract_text_from_image(self, file_link):
        from PIL import Image
        image = Image.open(file_link)
        text = pytesseract.image_to_string(image)
        return text

    def extract_text(self, file_link):
        pdf_text = ""
        with file_link.open() as f:
            pdf_stream = BytesIO(f.read())
            pdf_reader = PdfReader(pdf_stream)
            for page in pdf_reader.pages:
                pdf_text += page.extract_text().strip() + " "
        return pdf_text

    def file_summary_history(self, request):
        summary_history_objects = FileSumarizationModel.objects.filter(
            user_id=request.user.id,
            sub_category=5
        ).order_by("-created_at")
        pagination_obj = CustomPagination()
        search_keys = []
        result = pagination_obj.custom_pagination(
            request, search_keys, categorySerializer.GetFileSumarizationSerializer, summary_history_objects)
        return {"data": result, "message": messages.FETCH, "status": 200}

    def get_file_summary_by_id(self, request, file_id):
        try:
            get_summary_obj = FileSumarizationModel.objects.get(id=file_id)
        except FileSumarizationModel.DoesNotExist:
            return {"data": None, "message": messages.RECORD_NOT_FOUND, "status": 400}
        serializer = categorySerializer.GetFileSummarizationIdSerializer(
            get_summary_obj)
        return {"data": serializer.data, "message": messages.SUMMARY_FETCHED, "status": 200}

    def generate_file_important_vocabulary(self, request):
        # Get text_data from the request payload
        from decouple import config
        import openai
        openai.api_key = config("OPENAI_KEY")
        API_KEY = config("OPENAI_KEY")
        text_data = request.data.get('text_data')

        try:
            message=[
                        {"role": "system", "content": "You are a helpful assistant designed to output JSON."},
                        {"role": "user", "content": [
                                {"type": "text", "text": "Generate only important vocabulary words you found in the input which you consider to be important as a reference point. Format should be python list of words. If there are no important vocabulary just return 'Cannot find vocabulary'"},
                                {"type": "text", "text": f"Input data: {text_data}"}
                            ]}
                    ]
            chatbot = openai.ChatCompletion.create(
                model="gpt-3.5-turbo", messages=message,response_format={ "type": "json_object" },temperature = 0.0,
            )
            reply = chatbot.choices[0].message.content
            result = json.loads(reply)
            if isinstance(list(result.values())[0], str):
                final_response = []
            else:
               final_response = list(result.values())[0]
            return {"data": final_response, "message": "Vocabulary generated successfully", "status": 200}
        except Exception as e:
            print(str(e), type(e), '-----ee--e-e-e-e-e-e--e-----')
            return {"error": str(e), "message": messages.WENT_WRONG, "status": 400}

    # def extract_text(self, file_link):
    #     pdf_text = ""
    #     with file_link.open() as f:
    #         pdf_stream = BytesIO(f.read())
    #         pdf_reader = PdfReader(pdf_stream)
    #         for page in pdf_reader.pages:
    #             pdf_text += page.extract_text().strip() + "\n"
    #     return pdf_text
    ############# FILE CONVERSIONS ###############

    # def word_to_pdf(self , request):
    #     """
    #         We have to provide the link of the doc file and after
    #         that we can convert the file into pdf file.
    #     """
    #     word_file = request.FILES.get("pdf_file")
    #     file_name = word_file.name

    #     doc = aw.Document(word_file)
    #     doc.save(f"{file_name}.pdf")
    #     return {"data": "", "message": "conversion is done", "status": 200}

    def pdf_to_word(self, request):
        try:
            pdf_file = request.FILES.get("pdf_file")
            try:
                pdf_file: UploadedFile
                doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
                if doc.is_encrypted:
                    return {"data": None, "message": messages.PROTECTED_PDF, "status": 400}
            except Exception as e:
                pass
            file_name = generate_file_name(pdf_file.name)[0]
            OUTPUT_FILE_NAME = file_name
            file_name = f"{file_name}_{random.randint(10000, 99999)}"
            output_word_file = f"{OUTPUT_FILE_NAME}.docx"
            input_name = f"{file_name}_{random.randint(10000, 99999)}"
            input_pdf_file = f"{input_name}.pdf"
            delete_files = [input_pdf_file, output_word_file]

            fs = FileSystemStorage()
            fs.save(input_pdf_file, pdf_file)
            cv = Converter(input_pdf_file)
            # cv = Converter(BytesIO(pdf_content))
            cv.convert(output_word_file, start=0, end=None)
            cv.close()
            SAVED_FILE_RESPONSE = save_file_conversion(
                output_word_file, output_word_file, "word")
            data = {
                "media_url": SAVED_FILE_RESPONSE[0],
                "media_type": "word",
                "media_name": SAVED_FILE_RESPONSE[1]
            }
            serializer = CreateUpdateUploadMediaSerializer(data=data)
            if serializer.is_valid():
                serializer.save()
            for file in delete_files:
                if os.path.exists(file):
                    os.remove(file)
            save_file_in_model = FileConversationModel.objects.create(
                user_id=request.user.id,
                converted_media_id=serializer.data["id"],
                sub_category=10
            )
            return {"data": data, "message": messages.PDF_TO_WORD, "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.PLEASE_UPLOAD_AGAIN, "status": 400}

    def convert_pdf_to_excel(self, request):
        excel_file = request.FILES.get("pdf_file")
        try:
            excel_file: UploadedFile
            doc = fitz.open(stream=excel_file.read(), filetype="pdf")
            if doc.is_encrypted:
                return {"data": None, "message": messages.PROTECTED_PDF, "status": 400}
        except Exception as e:
            pass
        file_name = f"output_{random.randint(10000, 99999)}"
        OUTPUT_FILE_NAME = generate_file_name(excel_file.name)[0]
        pdf_content = excel_file.read()

        # Convert the PDF content to a byte stream
        pdf_byte_stream = BytesIO(pdf_content)

        # excel_path = str(excel_file)
        output_path = f"{file_name}.xlsx"

        self.pdf_to_excel(excel_file, output_path)
        SAVED_FILE_RESPONSE = save_file_conversion(
            output_path, f"{OUTPUT_FILE_NAME}.xlsx", "excel")
        data = {
            "media_url": SAVED_FILE_RESPONSE[0],
            "media_type": "excel",
            "media_name": SAVED_FILE_RESPONSE[1]
        }
        serializer = CreateUpdateUploadMediaSerializer(data=data)
        if serializer.is_valid():
            serializer.save()
        if os.path.exists(output_path):
            os.remove(output_path)
        save_file_in_model = FileConversationModel.objects.create(
            user_id=request.user.id,
            converted_media_id=serializer.data["id"],
            sub_category=14
        )
        return {"data": data, "message": messages.PDF_TO_EXCEL, "status": 200}

    def word_to_pdf(self, request):
        word_file = request.FILES.get("word_file")
        if not word_file:
            return {"message": "No Word file provided", "status": 400}

        # Generate unique file names
        file_name = "".join((word_file.name).split(" "))
        base_name = f"output_{random.randint(10000, 99999)}"
        temp_dir = tempfile.gettempdir()
        input_word_file = os.path.join(temp_dir, f"{base_name}.docx")
        # output_pdf_file = os.path.join(temp_dir, f"{base_name}.pdf")
        output_pdf_file = os.path.join(f"{base_name}.pdf")

        # Save the uploaded Word file temporarily
        with open(input_word_file, 'wb') as f:
            for chunk in word_file.chunks():
                f.write(chunk)

        document = Document(input_word_file)
        content = []
        for paragraph in document.paragraphs:
            content.append(paragraph.text)

        # Generate PDF using ReportLab
        doc = SimpleDocTemplate(output_pdf_file, pagesize=letter)
        styles = getSampleStyleSheet()
        paragraphs = [Paragraph(text, styles["Normal"]) for text in content]
        doc.build(paragraphs)
        # Convert Word to PDF
        convert(input_word_file, output_pdf_file)
        # Handle the converted PDF
        SAVED_FILE_RESPONSE = save_file_conversion(
            output_pdf_file, output_pdf_file, "application/pdf")
        data = {
            "media_url": SAVED_FILE_RESPONSE[0],
            "media_type": "pdf",
            "media_name": SAVED_FILE_RESPONSE[1]
        }
        serializer = CreateUpdateUploadMediaSerializer(data=data)
        if serializer.is_valid():
            serializer.save()
        save_file_in_model = FileConversationModel.objects.create(
            user_id=request.user.id,
            converted_media_id=serializer.data["id"],
            sub_category=11
        )
        # Save a copy of the output PDF file to a designated directory on your system
        # designated_dir = os.path.join(settings.BASE_DIR, 'saved_pdf_files')
        # os.makedirs(designated_dir, exist_ok=True)
        # final_pdf_path = os.path.join(designated_dir, f"{base_name}.pdf")

        # Copy the PDF file to the designated directory
        # with open(final_pdf_path, 'wb') as final_pdf_file:
        #     with open(output_pdf_file, 'rb') as temp_pdf_file:
        #         final_pdf_file.write(temp_pdf_file.read())

        # Clean up temporary files
        if os.path.exists(input_word_file):
            os.remove(input_word_file)
        if os.path.exists(output_pdf_file):
            os.remove(output_pdf_file)

        return {
            "data": data,
            "message": "done",
            "status": 200
        }

    def pdf_to_excel(self, pdf_path, excel_path):
        tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)

        # Create an Excel writer object
        with pd.ExcelWriter(excel_path, engine='xlsxwriter') as writer:
            # Iterate through each table and write it to a separate sheet
            for i, table in enumerate(tables):
                # Write the table to the Excel sheet
                table.to_excel(writer, sheet_name=f"Sheet {i+1}", index=False)

                # Access the worksheet object
                worksheet = writer.sheets[f"Sheet {i+1}"]

                # Iterate through each column and set its width
                for idx, col in enumerate(table.columns):
                    series = table[col]
                    max_len = max((
                        series.astype(str).map(len).max(),
                        len(str(series.name))
                    )) + 1
                    worksheet.set_column(idx, idx, max_len)

    def excel_to_pdf(self, request):
        try:
            excel_file = request.FILES.get("excel_file")
            file_name = excel_file.name
            OUTPUT_FILE_NAME = generate_file_name(file_name)[0] + ".pdf"
            file_name = file_name.replace(" ", "")
            file_name = file_name.replace("%", "")

            # excel_path = str(excel_file)
            file_save_path = f"{file_name}_{random.randint(10000, 99999)}.pdf"

            self.Excel_To_Pdf(excel_file, file_save_path)
            SAVED_FILE_RESPONSE = save_file_conversion(
                file_save_path, OUTPUT_FILE_NAME, "application/pdf")
            data = {
                "media_url": SAVED_FILE_RESPONSE[0],
                "media_type": "excel",
                "media_name": SAVED_FILE_RESPONSE[1]
            }
            serializer = CreateUpdateUploadMediaSerializer(data=data)
            if serializer.is_valid():
                serializer.save()
            if os.path.exists(file_save_path):
                os.remove(file_save_path)
            save_file_in_model = FileConversationModel.objects.create(
                user_id=request.user.id,
                converted_media_id=serializer.data["id"],
                sub_category=15
            )
            return {"data": serializer.data, "message": messages.CONVERT_SUCCESS, "status": 200}
        except ValueError as ve:
            if "must have at least a row and column" in str(ve):
                return {"data": None, "message": messages.EMPTY_EXCEL, "status": 400}
            return {"data": str(ve), "message": messages.PLEASE_UPLOAD_AGAIN, "status": 400}
        except Exception as err:
            return {"data": str(err), "message": messages.PLEASE_UPLOAD_AGAIN, "status": 400}

    def Excel_To_Pdf(self, excel_file, pdf_file):
        # Read Excel file into pandas DataFrame
        df = pd.read_excel(excel_file)
        # Convert DataFrame to list of lists (2D array)
        data = [df.columns.tolist()] + df.values.tolist()
        # Create PDF
        doc = SimpleDocTemplate(pdf_file, pagesize=A3)
        table = Table(data)
        # Style the table
        style = TableStyle([('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                            ('GRID', (0, 0), (-1, -1), 1, colors.black)])

        table.setStyle(style)
        # Add table to PDF
        doc.build([table])

    # def convert_pdf_to_image(self, request):
    #     pdf_file = request.FILES.get("pdf_file")
    #     file_name = "".join(str(pdf_file).split(" "))
    #     # file_save_path= f"image_{random.randint(10000, 99999)}.jpg"
    #     file_name = f"{random.randint(10000, 99999)}_{file_name}"
    #     input_pdf_file = f"{random.randint(10000, 99999)}_{file_name}"
    #     # input_pdf_file = f"{input_name}"

    #     fs = FileSystemStorage()
    #     fs.save(input_pdf_file, pdf_file)
    #     image_path_prefix = file_name.replace("pdf", "jpg")
    #     saved_files = self.pdf_to_image(input_pdf_file, image_path_prefix)
    #     if os.path.exists(input_pdf_file):
    #         os.remove(input_pdf_file)
    #     return {"data": saved_files, "message": messages.CONVERT_SUCCESS, "status": 200}

    # def pdf_to_image(self, pdf_path, image_path_prefix):
    #     # Convert PDF to a list of PIL Image objects
    #     images = convert_from_path(pdf_path)
    #     images_data = []
    #     # Save each page as an image
    #     for i, image in enumerate(images):
    #         image_path = f"{random.randint(10000, 99999)}_{image_path_prefix}"  # Change the extension as needed
    #         image.save(image_path, 'JPEG')  # Change the format as needed
    #         SAVED_FILE_RESPONSE = save_file_conversion(image_path, image_path, "jpg")
    #         data = {
    #                     "media_url": SAVED_FILE_RESPONSE[0],
    #                     "media_type": "image",
    #                     "media_name": SAVED_FILE_RESPONSE[1]
    #                 }
    #         serializer = CreateUpdateUploadMediaSerializer(data = data)
    #         if serializer.is_valid():
    #             serializer.save()
    #             images_data.append(serializer.data)
    #         if os.path.exists(image_path):
    #             os.remove(image_path)
    #     return images_data

    def convert_pdf_to_image(self, request):
        pdf_file = request.FILES.get("pdf_file")
        OUTPUT_FILE_NAME = generate_file_name(pdf_file.name)[0] + ".jpg"
        # Generate a unique file save path
        base_name = f"output_{random.randint(10000, 99999)}"
        temp_dir = tempfile.gettempdir()
        pdf_path = os.path.join(f"{base_name}.pdf")
        image_base_path = os.path.join(base_name)

        # Save the uploaded PDF file temporarily
        with open(pdf_path, 'wb') as f:
            for chunk in pdf_file.chunks():
                f.write(chunk)

        # Convert PDF to images
        conversion_result = self.convert_pdf_to_image2(
            pdf_path, image_base_path)
        if conversion_result["status"] != 200:
            return conversion_result

        # Handle the converted images
        image_paths = conversion_result["image_paths"]
        saved_file_responses = []
        start = 0
        for img_path in image_paths:
            image_name = OUTPUT_FILE_NAME + "_" + f"{start}" + ".jpg"
            SAVED_FILE_RESPONSE = save_file_conversion(
                img_path, image_name, "image/png")
            saved_file_responses.append({
                "media_url": SAVED_FILE_RESPONSE[0],
                "media_type": "image",
                "media_name": SAVED_FILE_RESPONSE[1]
            })
            start += 1
        images_ids = []
        for data in saved_file_responses:
            serializer = CreateUpdateUploadMediaSerializer(data=data)
            if serializer.is_valid():
                img_obj = serializer.save()
                images_ids.append(img_obj.id)

        if os.path.exists(pdf_path):
            os.remove(pdf_path)
        for img_path in image_paths:
            if os.path.exists(img_path):
                os.remove(img_path)

        save_file_in_model = FileConversationModel.objects.create(
            user_id=request.user.id,
            images=images_ids,
            sub_category=12
        )
        return {
            "data": saved_file_responses,
            "message": messages.CONVERT_SUCCESS,
            "status": status.HTTP_200_OK
        }

    def convert_pdf_to_image2(self, pdf_path, image_base_path):
        try:
            # Open the PDF file
            pdf_document = fitz.open(pdf_path)
            image_paths = []

            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)

                # Convert PDF page to image
                pix = page.get_pixmap()
                img_path = os.path.join(
                    f"{image_base_path}_page_{page_num}.png")
                pix.save(img_path)
                image_paths.append(img_path)

            return {
                "message": "Conversion successful",
                "status": 200,
                "image_paths": image_paths
            }
        except Exception as e:
            return {
                "message": f"Conversion failed: {str(e)}",
                "status": 500,
                "image_paths": []
            }

    def image_to_pdf(self, request):
        from PIL import Image
        try:
            # Get a list of uploaded images
            image_files = request.FILES.getlist("image")
            if not image_files:
                return {"message": "No images provided", "status": 400}

            fs = FileSystemStorage()
            saved_image_paths = []

            # Generate a unique file name for the PDF
            # Use the first image's name as the base
            file_name = os.path.splitext(image_files[0].name)[0]
            output_pdf_file = f"{file_name}_{random.randint(10000, 99999)}.pdf"
            file_save_path = fs.path(output_pdf_file)

            # Create a PDF document
            c = canvas.Canvas(file_save_path, pagesize=letter)

            for image_file in image_files:
                f_name = image_file.name
                f_name = f_name.replace(" ", "")
                input_image_file = fs.save(f_name, image_file)
                saved_image_paths.append(fs.path(input_image_file))

                # Open the image file
                with Image.open(fs.path(input_image_file)) as img:
                    img_width, img_height = img.size

                    # Convert the image size to fit into the PDF page size
                    pdf_width, pdf_height = letter
                    aspect = img_width / img_height

                    if img_width > img_height:
                        # Landscape orientation
                        img_width = pdf_width
                        img_height = pdf_width / aspect
                    else:
                        # Portrait orientation
                        img_height = pdf_height
                        img_width = pdf_height * aspect

                    x = (pdf_width - img_width) / 2
                    y = (pdf_height - img_height) / 2

                    # Insert the image into the PDF
                    c.drawImage(fs.path(input_image_file), x, y,
                                width=img_width, height=img_height)
                    c.showPage()  # Add a page break

            # Save the PDF document
            c.save()

            # Generate the URL for the PDF file
            SAVED_FILE_RESPONSE = save_file_conversion(
                output_pdf_file, "output.pdf", "application/pdf")
            data = {
                "media_url": SAVED_FILE_RESPONSE[0],
                "media_type": "pdf",
                "media_name": SAVED_FILE_RESPONSE[1]
            }
            serializer = CreateUpdateUploadMediaSerializer(data=data)
            if serializer.is_valid():
                serializer.save()

            # Clean up saved image files
            for saved_image_path in saved_image_paths:
                if os.path.exists(saved_image_path):
                    os.remove(saved_image_path)
            if os.path.exists(file_save_path):
                os.remove(file_save_path)

            save_file_in_model = FileConversationModel.objects.create(
                user_id=request.user.id,
                converted_media_id=serializer.data["id"],
                sub_category=13
            )
            return {"data": serializer.data, "status": 200}

        except Exception as e:
            return {"message": str(e), "status": 400}

        #     file_name = image_file.name

        #     doc = aw.Document()
        #     builder = aw.DocumentBuilder(doc)

        #     builder.insert_image(str(image_file))

        #     doc.save(f"{file_name}.pdf")

        #     return {"message":messages.CONVERT_SUCCESS, "status": 200}
        # except Exception as e:
        #     return {"message": str(e), "status": 400}

    def ppt_to_pdf(self, request):
        ppt_file = request.FILES.get("ppt_file")
        file_name = ppt_file.name

        # Generate a unique file save path
        base_name = f"output_{random.randint(10000, 99999)}"
        temp_dir = tempfile.gettempdir()
        # ppt_path = os.path.join(temp_dir, f"{base_name}.pptx")
        ppt_path = f"{base_name}.pptx"
        pdf_path = os.path.join(f"{base_name}.pdf")

        # Save the uploaded PPT file temporarily
        with open(ppt_path, 'wb') as f:
            for chunk in ppt_file.chunks():
                f.write(chunk)

        # Convert PPT to PDF using images
        conversion_result = self.convert_ppt_to_pdf_with_images(
            ppt_path, pdf_path)

        if conversion_result["status"] != 200:
            return conversion_result

        # Handle the converted PDF
        SAVED_FILE_RESPONSE = save_file_conversion(
            pdf_path, pdf_path, "application/pdf")
        data = {
            "media_url": SAVED_FILE_RESPONSE[0],
            "media_type": "ppt",
            "media_name": SAVED_FILE_RESPONSE[1]
        }
        serializer = CreateUpdateUploadMediaSerializer(data=data)
        if serializer.is_valid():
            serializer.save()
        save_file_in_model = FileConversationModel.objects.create(
            user_id=request.user.id,
            converted_media_id=serializer.data["id"],
            sub_category=17
        )
        # Save a copy of the output PDF file to a designated directory on your system
        # designated_dir = os.path.join(tempfile.gettempdir(), 'saved_ppt_pdf_files')
        # os.makedirs(designated_dir, exist_ok=True)
        # final_pdf_path = os.path.join(designated_dir, f"{base_name}.pdf")
        # with open(final_pdf_path, 'wb') as final_pdf_file:
        #     with open(pdf_path, 'rb') as temp_pdf_file:
        #         final_pdf_file.write(temp_pdf_file.read())

        # Clean up temporary files
        if os.path.exists(ppt_path):
            os.remove(ppt_path)
        if os.path.exists(pdf_path):
            os.remove(pdf_path)

        return {
            "data": serializer.data,
            "message": messages.CONVERT_SUCCESS,
            "status": status.HTTP_200_OK
        }

    def convert_ppt_to_pdf_with_images(self, ppt_path, pdf_path):
        try:
            prs = Presentation(ppt_path)
            temp_dir = tempfile.gettempdir()
            slide_images = []

            for slide in prs.slides:
                # slide_img_base = os.path.join(temp_dir, f"slide_{prs.slides.index(slide)}")
                slide_path = f"slide_{prs.slides.index(slide)}"
                success = self.save_slide_as_image(slide, slide_path)
                if success:
                    slide_images.append(slide_img_base)
                else:
                    return {"message": "Error while converting.", "status": 400}
                # return {"message": "slide saved", "status": 400}
            if slide_images:
                # Assuming the first image's index is 0
                first_image = Image.open(slide_images[0] + "_0.jpg")
                first_image.save(pdf_path, save_all=True, append_images=[
                                 Image.open(f"{img}_0.png") for img in slide_images[1:]])

            # Clean up temporary image files
            for slide_img_base in slide_images:
                for idx in range(len(slide.shapes)):
                    os.remove(f"{slide_img_base}_{idx}.png")

            return {"message": "Conversion successful", "status": 200}
        except Exception as e:
            return {"message": f"Conversion failed: {str(e)}", "status": 500}

    def save_slide_as_image(self, slide, output_path):
        """
        Save a PowerPoint slide as an image.

        Args:
            slide (pptx.slide.Slide): The slide to save as an image.
            output_path (str): The file path to save the image.
            width (int, optional): The width of the output image. Defaults to 1024.
            height (int, optional): The height of the output image. Defaults to 768.
        """
        try:
            img = Image.new(mode="RGB", size=(1024, 768))
            img_draw = ImageDraw.Draw(img)
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img.paste(shape.image, (shape.left, shape.top))
                else:
                    img_draw.text((shape.left, shape.top),
                                  shape.text, fill='white')
                img.save(f"{output_path}.jpg")
            return True
        except Exception as err:
            return False

    # def save_slide_as_image(self, slide, img_path_base):
    #     try:
    #         # os.makedirs(os.path.dirname(img_path_base), exist_ok=True)
    #         for idx, shape in enumerate(slide.shapes):
    #             if not hasattr(shape, 'image'):
    #                 continue
    #             image = shape.image
    #             image_bytes = image.blob
    #             img_path = f"{img_path_base}_{idx}.png"
    #             with open(img_path, 'wb') as img_file:
    #                 img_file.write(image_bytes)
    #         return True  # Indicate successful operation
    #     except Exception as e:
    #         return False  # Indicate failure

    def pdf_to_ppt(self, request):
        pdf_file = request.FILES.get("pdf_file")
        OUTPUT_FILE_NAME = generate_file_name(pdf_file.name)[0] + ".pptx"
        # Generate a unique file save path
        base_name = f"output_{random.randint(10000, 99999)}"
        temp_dir = tempfile.gettempdir()
        pdf_path = os.path.join(temp_dir, f"{base_name}.pdf")
        ppt_path = os.path.join(f"{base_name}.pptx")

        # Save the uploaded PDF file temporarily
        with open(pdf_path, 'wb') as f:
            for chunk in pdf_file.chunks():
                f.write(chunk)

        # Convert PDF to PPT using the instance method
        conversion_result = self.convert_pdf_to_ppt(pdf_path, ppt_path)

        if conversion_result["status"] != 200:
            return conversion_result

        # Handle the converted PPT (example save_file_conversion function call)
        SAVED_FILE_RESPONSE = save_file_conversion(
            ppt_path, OUTPUT_FILE_NAME, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        data = {
            "media_url": SAVED_FILE_RESPONSE[0],
            "media_type": "pptx",
            "media_name": SAVED_FILE_RESPONSE[1]
        }

        # Example serializer usage (replace with your serializer)
        serializer = CreateUpdateUploadMediaSerializer(data=data)
        if serializer.is_valid():
            serializer.save()

        # Example model creation (replace with your model)
        save_file_in_model = FileConversationModel.objects.create(
            user_id=request.user.id,
            converted_media_id=serializer.data["id"],
            sub_category=16
        )

        # Clean up temporary files
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
        if os.path.exists(ppt_path):
            os.remove(ppt_path)

        return {
            "data": serializer.data,
            "message": messages.CONVERT_SUCCESS,
            "status": status.HTTP_200_OK
        }

    def convert_pdf_to_ppt(self, pdf_path, ppt_path):
        from pptx.util import Pt
        from PIL import Image
        from pptx import Presentation
        try:
            # Open the PDF file
            pdf_document = fitz.open(pdf_path)
            temp_dir = tempfile.gettempdir()

            prs = Presentation()

            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)

                # Convert PDF page to image
                pix = page.get_pixmap()
                img_path = os.path.join(temp_dir, f"page_{page_num}.png")
                pix.save(img_path)

                # Open the image and get its size
                with Image.open(img_path) as img:
                    img_width, img_height = img.size

                    # Define slide dimensions for portrait orientation (switching width and height)
                    # Height in inches (adjust as needed)
                    slide_height = Inches(10)
                    # Width in inches (adjust as needed)
                    slide_width = Inches(7.5)

                    # Calculate scaling factor to fit image within slide dimensions
                    scale = min(slide_width / img_width,
                                slide_height / img_height)

                    # Calculate centered position within slide
                    left = (slide_width - img_width * scale) / 2
                    top = (slide_height - img_height * scale) / 2

                    # Add a blank slide with portrait orientation
                    # Index 6 for blank slide layout (Title Slide)
                    slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(slide_layout)

                    # Set slide dimensions and orientation
                    slide.slide_width = Pt(slide_width.pt)
                    slide.slide_height = Pt(slide_height.pt)
                    prs.slide_width = slide.slide_width
                    prs.slide_height = slide.slide_height

                    # Add image to slide
                    slide.shapes.add_picture(
                        img_path, left, top, width=img_width * scale, height=img_height * scale)

                # Remove the temporary image file
                os.remove(img_path)

            # Save the presentation
            prs.save(ppt_path)

            return {"message": "Conversion successful", "status": 200}
        except Exception as e:
            return {"message": f"Conversion failed: {str(e)}", "status": 500}

    def file_conversions_history(self, request):
        try:
            all_objs = FileConversationModel.objects.filter(
                user=request.user.id).order_by("-created_at")
            pagination_obj = CustomPagination()
            search_keys = []
            result = pagination_obj.custom_pagination(
                request, search_keys, categorySerializer.FileConversionlistingSerializer, all_objs)
            return {"data": result, "message": messages.FETCH, "status": 200}
        except:
            return {"data": None, "message": "Something went wrong", "status": 400}


# note

    def voice_to_text(self, request):
        try:
            audio = request.FILES.get("audio")
            recognizer = sr.Recognizer()
            audio = AudioSegment.from_mp3(audio)
            audio.export("audio.wav", format="wav")
            with sr.AudioFile("audio.wav") as source:
                audio = recognizer.record(source)
            text = recognizer.recognize_google(audio, language='en-US')
            return {"message": text, "status": 200}
        except Exception as e:
            return {"message": str(e), "status": 400}

    def add_notes(self, request):
        try:
            serializers = categorySerializer.AddNoteSerializer(
                data=request.data)
            if serializers.is_valid():
                user_obj = serializers.save()
                user_obj.user_id = request.user.id
                user_obj.save()
            return {"data": serializers.data, "message": "note uploaded successfully", "status": 200}

        except Exception as e:
            return {"error": str(e), "message": messages.WENT_WRONG, "status": 400}

    def ai_explanation(self, request):
        text = request.data.get("text")
        llm = ChatGoogleGenerativeAI(model="gemini-pro")
        try:
            response = llm.invoke(text)
            result_qu = to_markdown(response.content)
            return {"data": result_qu, "message": messages.FETCH, "status": 200}
        except Exception as e:
            return {"error": str(e), "message": messages.WENT_WRONG, "status": 400}

    def change_language_note(self, request):
        try:
            text = request.data.get("text")
            translator = GoogleTranslator(source="auto", target="ar")
            answer = translator.translate(text)
            try:
                answer = ast.literal_eval(answer)
            except Exception as err:
                pass
            return {"data": answer, "message": messages.FETCH, "status": 200}
        except Exception as e:
            return {"error": str(e), "message": messages.WENT_WRONG, "status": 400}

    def get_all_listing_notes(self, request):
        try:
            notes_obj = NoteModel.objects.filter(user_id=request.user.id)
            serializer = categorySerializer.GetNoteSerializer(
                notes_obj, many=True)
            return {"data": serializer.data, "message": messages.FETCH, "status": 200}
        except:
            return {"data": None, "message": messages.WENT_WRONG, "status": 400}

    def get_notes_by_id(self, request, id):
        try:
            notes_obj = NoteModel.objects.filter(user_id=request.user.id)
            serializer = categorySerializer.GetNoteSerializer(notes_obj)
            return {"data": serializer.data, "message": messages.FETCH, "status": 200}
        except:
            return {"data": None, "message": messages.WENT_WRONG, "status": 400}


# research

    def to_markdown(text):
        text = text.replace('*', '')
        intent_text = (textwrap.indent(text, '', predicate=lambda _: True))
        return intent_text

    def get_research_answer(self, request):
        reduce_citation = request.data.get("reduce_citation")
        description = request.data.get("description")
        if not request.data.get('upload_reference'):
            topic = request.data.get("topic")
            page = request.data.get("page")
            # words=int(page)*300
            tone = request.data.get("tone")
            reference = request.data.get("reference")
            data = f"You are a topics list generator. Generate research topics list based on {topic}. Output should contain topics headings(strictly numbered like 1,2,3,.....) and slide headings(strictly numbered like i, ii, iii , ......)."
            # data=f"You are a topics list generator. Generate research topics list based on {topic}. Output should contain only three topics headings(numbered like 1,2,3) and strictly two side headings(numbered like i, ii, iii)."
            query = data
            llm = ChatGoogleGenerativeAI(model="gemini-pro")
            try:
                response = llm.invoke(query)
                result = to_markdown(response.content)
                if "Invalid input provided" in result:
                    return {"data": "Invalid input provided", "message": "Invalid input provided", "status": 400}
                save_to_db = CategoryModel.objects.create(
                    user_id=request.user.id,
                    topic=topic,
                    page=page,
                    tone=tone,
                    reference=reference,
                    category=4,
                    result=result
                )
                return {"data": result, "record_id": save_to_db.id, "message": messages.FETCH, "status": 200}
            except Exception as e:
                print(e, '------ererro---------')
                return {"data": str(e), "message": messages.WENT_WRONG, "status": 400}

        pdf_link = request.data['upload_reference']
        data = f"generate esaay of mininimum 500 words from given link {pdf_link} and it should define with this {description} and reduce citation will be {reduce_citation}"
        query = data
        llm = ChatGoogleGenerativeAI(model="gemini-pro")
        try:
            response = llm.invoke(query)
            result = to_markdown(response.content)
            return {"data": result, "message": messages.FETCH, "status": 200}
        except Exception as e:
            return {"data": str(e), "message": messages.WENT_WRONG, "status": 400}

    def regenerate_research_solution(self, request, id):
        PAGE_REFERENCES = {1: (600, 1200), 2: (
            1800, 3000), 3: (4200, 7200), 4: (9000, 15000)}
        get_research_record = CategoryModel.objects.get(id=id)
        llm = ChatGoogleGenerativeAI(model="gemini-pro")
        if request.GET.get("page") in ["1/", 1, "1"]:
            try:
                if get_research_record.research_type == 1:
                    topic = get_research_record.topic
                    # min_pages = PAGE_REFERENCES[get_research_record.page][0]
                    # max_pages = PAGE_REFERENCES[get_research_record.page][1]
                    # tone = get_research_record.tone
                    # reference = get_research_record.reference
                    # reduce_citations = get_research_record.reduced_citations
                    # tone = get_research_record.tone
                    # QUERY=F"You are research generator. Generate some theory on the topics which I provide you. Whole research should be of approximately {min_pages} to {max_pages} words with voice of tone as {tone} and take reference from {reference} with reduce citations as {reduce_citations}. Format should be descriptive. Keep the same name as topic heading which I provide you and also in same order of topics. Strictly keep Heading(numbered as 1,2,3) and side headings(numbered as i, ii, iii)."
                    QUERY = f"You are a topics list generator. Generate research topics list based on {topic}. Output should contain only three topics headings(numbered like 1,2,3) and strictly two side headings(numbered like i, ii, iii)."
                    response = llm.invoke(QUERY)

                    result = to_markdown(response.content)
                    return {"data": result, "message": "Research topics generated successfully", "status": 200}
                elif get_research_record.research_type == 2:
                    image_links = get_research_record.research_file_links
                    QUERY = f"You are topics list generator. Generate research topics list based on links I provide to you. Output should contain only three topics headings(numbered like 1,2,3) and strictly two side headings(numbered like i, ii, iii)."
                    # message_content = [
                    #     {
                    #         "type": "text",
                    #         "text": QUERY,
                    #     }
                    # ]
                    # for image_link in image_links:
                    #     message_content.append({
                    #         "type": "image_url",
                    #         "image_url": str(image_link)
                    #     })
                    response = pdf_processing(image_links[0], QUERY)
                    final_response = response.replace("*", "").replace("-", "")
                    # message = HumanMessage(content=message_content)
                    # response = llm.invoke([message])
                    # result = to_markdown(response.content)
                    return {"data": final_response, "message": "Research topics generated successfully", "status": 200}
            except Exception as err:
                return {"data": str(err), "message": messages.WENT_WRONG, "status": 400}
        else:
            result = self.generate_detailed_research_based_on_topics(
                request, id)
            return result

    def research_based_on_reference(self, request):
        try:
            llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash")
            description = request.data.get("description", "")
            reduce_citation = True if request.data.get(
                "reduce_citation") == "true" else False
            image_links = []
            for img in dict(request.data)["files"]:
                get_link = save_image(img)
                image_links.append(get_link[0])
            query = f"You are topics list generator. Generate research topics list based on links I provide to you with reduce citations as {reduce_citation}. Output should contain only three topics headings(numbered like 1,2,3) and strictly two side headings(numbered like i, ii, iii)."
            response = pdf_processing(image_links[0], query)
            if isinstance(response, dict):
                return {"data": None, "message": response["message"], "status": 400}
            # message_content = [
            #     {
            #         "type": "text",
            #         "text": query,
            #     }
            # ]
            # for image_link in image_links:
            #     message_content.append({
            #         "type": "image_url",
            #         "image_url": str(image_link)
            #     })
            # message = HumanMessage(content=message_content)
            # response = llm.invoke([message])
            final_response = response.replace("*", "").replace("-", "")
            save_to_db = CategoryModel.objects.create(
                user_id=request.user.id,
                description=description,
                category=4,
                research_type=2,
                reduced_citations=reduce_citation,
                result=final_response,
                research_file_links=image_links
            )
            return {"data": final_response, "record_id": save_to_db.id, "message": messages.RESEARCH_GENERATED, "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.WENT_WRONG, "status": 400}

    def generate_detailed_research_based_on_topics(self, request, id):
        PAGE_REFERENCES = {1: (600, 1200), 2: (
            1800, 3000), 3: (4200, 7200), 4: (9000, 15000)}
        get_research_record = CategoryModel.objects.get(id=id)
        try:
            api_type = 1
            min_pages = PAGE_REFERENCES[get_research_record.page][0]
            max_pages = PAGE_REFERENCES[get_research_record.page][1]
            tone = get_research_record.tone
            reference = get_research_record.reference
            reduce_citations = get_research_record.reduced_citations
        except KeyError:
            api_type = 2
            image_links = get_research_record.research_file_links
            descriptoin = get_research_record.description
            reduce_citations = get_research_record.reduced_citations
        except Exception as err:
            return {"data": str(err), "message": messages.WENT_WRONG, "status": 400}
        llm = ChatGoogleGenerativeAI(model="gemini-pro")
        html_text = request.data.get("html_text")
        # extract text from html
        if html_text:
            soup = BeautifulSoup(html_text, "html.parser")
            all_topics = []
            for ele in soup.find_all(["h1", "li"]):
                all_topics.append(ele.get_text())
                get_research_record.all_topics = all_topics
                get_research_record.save()
        else:
            all_topics = get_research_record.all_topics
        ####
        if api_type == 1:
            QUERY = f"You are research generator. Generate some theory on the topics which I provide you. Whole research should be of approximately {min_pages} to {max_pages} words with voice of tone as {tone} and take reference from {reference} with reduce citations as {reduce_citations}. Format should be descriptive. Keep the same name as topic heading which I provide you and also in same order of topics. Strictly keep Heading(numbered as 1,2,3) and side headings(numbered as i, ii, iii)."
        elif api_type == 2:
            QUERY = f"You are research generator. Generate some theory on the topics which I provide you. Whole research should be of approximately 400 to 800 words along with reduce citations as {reduce_citations}. Format should be descriptive. Keep the same name as topic heading which I provide you in list and also in same order of topics. Strictly keep Headings (strictly numbered as 1,2,3) and side headings(strictly numbered as i, ii, iii)."
        message_content = [
            {
                "type": "text",
                "text": QUERY
            }
        ]
        for top in all_topics:
            message_content.append({
                "type": "text",
                "text": top
            })
        message = HumanMessage(content=message_content)
        response = llm.invoke([message])
        final_response = response.content.replace(
            "*", "").replace("#", "").replace("-", "")
        # save record
        # get_research_record.result = final_response
        # get_research_record.save()
        ##
        print(final_response, '-------final_response------')
        return {"data": final_response, "message": "Detailed research generated successfully", "status": 200}

    def save_research_topic_list(self, request, id):
        research_obj = CategoryModel.objects.get(id=id)
        research_obj.result = request.data.get("text")
        research_obj.save()
        return {"data": "", "message": messages.RESEARCH_SAVED, "status": 200}

    def save_rsearch_file(self, request):
        try:
            save_pdf = CategoryModel.objects.create(
                user_id=request.user.id, category=4, sub_category=5, media_id=request.data.get("pdf_file"))
            save_pdf.save()
            return {"data": request.data, "message": "saved successfully", "status": 200}
        except Exception as e:
            return {"data": str(e), "message": messages.WENT_WRONG, "status": 400}

    def get_history_research(self, request):
        try:
            research_obj = CategoryModel.objects.filter(
                user_id=request.user.id, category=4).order_by("-id")
            pagination_obj = CustomPagination()
            search_keys = []
            result = pagination_obj.custom_pagination(
                request, search_keys, categorySerializer.GetNoteListSerializer, research_obj)
            return {"data": result, "message": messages.FETCH, "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.WENT_WRONG, "status": 400}

    def get_research_by_id(self, request, id):
        try:
            research_obj = CategoryModel.objects.get(id=id)
            serializer = categorySerializer.GetNoteListSerializer(research_obj)
            return {"data": serializer.data, "message": messages.FETCH, "status": 200}
        except:
            return {"data": None, "message": messages.WENT_WRONG, "status": 400}

    def download_research_file(self, request, id):
        try:
            if not request.data.get("html_text").strip():
                return {"data": None, "message": "You cannot download empty file", "status": 400}
            assignment = CategoryModel.objects.get(id=id)
            if request.data["type"] == 1:
                if request.data["new"] is True:
                    file = self.html_to_pdf(request)
                    return {"data": file, "message": "pdf generated successfully", "status": 200}
                if assignment.download_file:
                    return {"data": assignment.download_file, "message": messages.UPDATED, "status": 200}
                file = self.html_to_pdf(request)
                assignment.download_file = file
                assignment.save()
            if request.data["type"] == 2:
                if request.data["new"] is True:
                    file = self.html_to_doc(request)
                    return {"data": file, "message": "pdf generated successfully", "status": 200}
                if assignment.download_doc_file:
                    return {"data": assignment.download_doc_file, "message": messages.UPDATED, "status": 200}
                file = self.html_to_doc(request)
                assignment.download_doc_file = file
                assignment.save()
            elif request.data["type"] == 3:
                file = self.html_to_pdf(request)
                Thread(target=send_pdf_file_to_mail, args=(
                    assignment.user.email, file)).start()
                return {"data": None, "message": "File send to your email successfully", "status": 200}
            return {"data": file, "message": messages.UPDATED, "status": 200}
        except CategoryModel.DoesNotExist:
            return {"data": None, "message": "Record not found", "status": 400}
        except Exception as e:
            return {"data": str(e), "message": messages.WENT_WRONG, "status": 400}

    def gemini_solution(self, file_link):
        llm = ChatGoogleGenerativeAI(model="gemini-pro")
        text_data = self.extract_text(file_link)
        message = HumanMessage(
            content=[
                {"type": "text",
                    # existing
                    # "text": "I am an invligator to mark the questions i need correct answers ,provide me correct answers for these questions and when needed diagrams and figures or explanations just give concise answers and give answers to remaining questions,lastly provide the answers in json list format (question no. ,question, options(this field will will only be there if options are present else no need ), correct answer),But if somehow you dont get proper pdf then fetch the words and give answer regarding words and dont give random response.but if its an scan image then extract text from it and provide solution regarding it"},
                    ## old ###
                    # "text": "I am an invligator i will give you scaned pdf of images. Whatever you find text in images give question and answer regarding them.And dont give random answers"},
                    ### new ####
                    "text": "You are a teacher. Generate questions and answers based on the data I provide to you. Format should be proper Python Javascript object notation list of dictionaries where every dictionary contains keys as 'question_no', 'question', 'correct_answer' and 'options'(if available)."},
                # "text": "I am an invligator to mark the questions i need correct answers ,provide me correct answers for these questions and when needed diagrams and figures or explanations just give concise answers and give answers to remaining questions,lastly provide the answers in json list format (question no. ,question, options(this field will will only be there if options are present else no need ), correct answer)"},
                {"type": "text", "text": text_data}
            ]
        )
        response = llm.invoke([message])
        result = to_markdown(response.content)
        return result

    def gemini_solution_review(self, file_link):
        llm = ChatGoogleGenerativeAI(model="gemini-pro")
        text_data = self.extract_text(file_link)
        message = HumanMessage(
            content=[
                {"type": "text",
                    "text": "These are some question and answers you have generated for my assigment previously.Now just add another key explanation and give explanation of that answer. And make sure to keep that response as it is by just adding explanation key, lastly provide the answers in json list of objects where each object should have keys (question_no, question, correct_answer, explanation, options(this field will will only be there if options are present else no need)"},
                #  "text": f"list the answers for all questions present  in these given file's (don't leave any question ,even if there is breaks between questions)and provide in  json  format (questtions which have no options just give correct answers in concise manner) try writing answer in this way  (question no. ,question, options(this field will will only be there if options are present else no need ),correct answer) "},
                {"type": "text", "text": text_data}
            ]
        )
        response = llm.invoke([message])
        result = to_markdown(response.content)
        return result

# assignment solution
    def text_translation(self, request):
        text = request.data.get("text")
        print(text, '--------')
        if isinstance(text, list):
            # query = "You are english to arabic translator. Translate all the words to arabic wherever you find which I provide you and don't translate the key names. Format should be python json list."
            query = f"""
                        First find the language of input and Translate to arabic if it is english or translate to english if it is arabic:
                        {text}. Strictly follow the format. Translate every question, every option and every answer.
                    """
            text = json.dumps(text)
            try:
                final_response = self.change_language_chatgpt(query, text)[
                    "questions"]
            except:
                final_response = self.change_language_chatgpt(query, text)

            # result = self.gemini_solution_for_text_translation(text, query)
            # try:
            #     final_response = json.loads(result)
            # except json.decoder.JSONDecodeError:
            #     final_response = ast.literal_eval(result)
            # except Exception as err:
            #     return {"data": None, "message": "Please try again", "status": 400}
            # print(final_response, '-----------final------')
        else:
            query = "You are english to arabic translator. Translate the text to arabic which I provide you.Output format should be proper human readable text ."
            result = self.gemini_solution_for_text_translation(text, query)
            final_response = result
        if isinstance(final_response, list):
            pass
        elif isinstance(final_response, dict):
            final_response = list(final_response.values())[0]
        return {"data": final_response, "message": "Text translated successfully.", "status": 200}

    def change_language_chatgpt(self, query, input_data):
        from decouple import config
        import openai
        openai.api_key = config("OPENAI_KEY")
        try:
            messages = [
                {"role": "system",
                    "content": "You are a helpful assistant designed to output JSON."},
                {"role": "user", "content": [
                    {"type": "text", "text": query},
                    {"type": "text", "text": input_data},
                ]}
            ]
            chatbot = openai.ChatCompletion.create(
                model="gpt-4o", messages=messages, response_format={"type": "json_object"}, temperature=0.0,
            )
            reply = chatbot.choices[0].message.content
            final_data = json.loads(reply)
            return final_data
        except:
            return []

    def gemini_solution_for_text_translation(self, text, query):
        llm = ChatGoogleGenerativeAI(model="gemini-pro")
        text_data = text
        message = HumanMessage(
            content=[
                {"type": "text",
                    "text": query},
                {"type": "text", "text": text_data}
            ]
        )
        response = llm.invoke([message])
        result = to_markdown(response.content)
        return result

    def get_assignment_solution(self, request):
        if int(request.data["type"]) == 1:
            file_link = request.FILES.get("file_link")
            try:
                final_response = []
                language = request.data.get("language", "english")
                # if request.data.get("language") == "arabic":
                print(language, '---language----')
                text_data = assignment_extract_text(file_link)
                for i in text_data:
                    result = assigment_chatGPT_pdf_processing(i, language)
                    print(result, '-----result----result-----')
                    if result:
                        if result.get("questions"):
                            final_response += result["questions"]
                        elif not result.get("questions") and isinstance(result, dict):
                            final_response.append(result)
                try:
                    for j, i in enumerate(final_response):
                        i["question_no"] = j + 1
                        if not i.get("options"):
                            i["question_type"] = 1
                        elif i["options"]:
                            i["question_type"] = 2
                except:
                    pass
                    # image_info = upload_media_obj.upload_media(request)
                final_data = AssignmentModel.objects.create(
                    user_id=request.user.id,
                    result=final_response
                )
                final_data.save()
                if not final_response:
                    return {"data": None, "message": "Please upload the file again", "status": 200}
                return {"data": final_response, "record_id": final_data.id, "message": "Assignment solution generated successfully", "status": 200}
            except Exception as e:
                print(e, '----eeee----eeee-----')
                return {"data": str(e), "message": "Please upload the file again", "status": 400}

        if int(request.data["type"]) == 2:
            try:
                # llm = ChatGoogleGenerativeAI(model="gemini-pro")
                images = dict(request.data)["file_link"]
                query = "You are a teacher. Generate questions and answers based on the data I provide to you. Format should be proper Python Javascript object notation list of dictionaries where every dictionary contains keys as 'question_no', 'question', 'correct_answer' and 'options'(if available)."
                gemini_result = []
                for file_image in images:
                    img = save_image(file_image)
                    result = self.image_processing_assignment_solution(
                        img[0], query)
                    # text_data = self.extract_text_from_image(file_image)
                    # message = HumanMessage(
                    #     content=[
                    #         {"type": "text",
                    #             "text": "You are a teacher. Generate questions and answers based on the data I provide to you. Format should be proper Python Javascript object notation list of dictionaries where every dictionary contains keys as 'question', 'answer' and 'options'(if available)."},
                    #         {"type": "text", "text":text_data}
                    #     ]
                    # )
                    # # result = self.gemini_solution(file_link)
                    # response = llm.invoke([message])
                    # result = to_markdown(response.content)
                    temp = self.format_final_response(result)
                    gemini_result += temp
                final_data = AssignmentModel.objects.create(
                    user_id=request.user.id,
                    result=gemini_result
                )
                final_data.save()
                return {"data": gemini_result, "record_id": final_data.id, "message": "Assignment solution generated successfully", "status": 200}
            except Exception as err:
                return {"data": str(err), "message": messages.PLEASE_UPLOAD_AGAIN, "status": 400}

    def image_processing_assignment_solution(self, image_link, query):
        '''processing the image and generate the mcq with options and answer 
            image_link is the s3 bucket link of image and query is string 
            which we use to generate the mcq of flashcards'''
        llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash")
        # example
        message = HumanMessage(
            content=[
                {
                    "type": "text",
                    "text": query,
                },
                {"type": "image_url", "image_url": str(image_link)},
            ]
        )
        response = llm.invoke([message])
        result_qu = to_markdown(response.content)
        return result_qu

    def format_final_response(self, result):
        final_response = ""
        try:
            for i in range(len(result)-1, -1, -1):
                if result[i] == "}":
                    break
            final_response = result[result.index("["): i+1] + "]"
            final_response = json.loads(final_response)
        except:
            pass
        try:
            for i in final_response:
                if not i.get("options"):
                    i["question_type"] = 1
                elif not i["options"]:
                    i["question_type"] = 1
                elif i["options"]:
                    i["question_type"] = 2
        except:
            pass
        return final_response

    # def get_all_assignment(self, request):
    #     try:
    #         data = AssignmentModel.objects.all()
    #         pagination_obj = CustomPagination()
    #         search_keys = []
    #         result = pagination_obj.custom_pagination(request, search_keys, categorySerializer.CreateAssignmentSerializers, data)
    #         return {"data":result,"message":messages.FETCH,"status":200}
    #     except Exception as e:
    #         return {"data":None,"message":messages.WENT_WRONG,"status":400}

    def get_assignment_solution_review(self, request, id):
        # file_link = request.FILES.get("file_link")
        file = self.get_assignment_solution_review_func(request)
        try:
            with open(file, 'rb') as f:
                file_content = f.read()
                in_memory_file = BytesIO(file_content)

        # Wrap the BytesIO object with Django's File class
            django_file = File(in_memory_file, name=os.path.basename(file))
            result = self.gemini_solution_review(django_file)
            final_response = ""
            last_ele = 0
            try:
                for i in range(len(result)-1, -1, -1):
                    if result[i] == "}":
                        last_ele = i
                        break
                final_response = result[result.index("["): i+1] + "]"
                final_response = json.loads(final_response)
            except ValueError:
                if result[0] == "{":
                    result = "[" + result[result.index("{"): last_ele+1] + "]"
                    final_response = json.loads(result)
            except Exception as err:
                return {"data": "Please try again.", "message": "Please try again", "status": 400}
            try:
                for i in final_response:
                    if not i.get("options"):
                        i["question_type"] = 1
                    elif not i["options"]:
                        i["question_type"] = 1
                    elif i["options"]:
                        i["question_type"] = 2
            except:
                pass
            try:
                for i in final_response:
                    # if i["explanation"] is not None:
                    #     i["explanation"]=i["correct_answer"]
                    if i["explanation"]:
                        i["correct_answer"] = i["explanation"]
            except Exception as err:
                pass
            if os.path.exists(file):
                os.remove(file)
            if not final_response:
                return {"data": "Empty Response", "message": "Please try again", "status": 400}
            return {"data": final_response, "message": "Review generated successfully", "status": 200}
        except Exception as error:
            return {"data": str(error), "message": "Please try again", "status": 400}

    def get_assignment_solution_review_func(self, request):
        html_text = request.data["html_text"]
        # Update this path as necessary
        path_to_wkhtmltopdf = 'C:\\Program Files\\wkhtmltopdf\\bin\\wkhtmltopdf.exe'
        config = pdfkit.configuration(wkhtmltopdf=path_to_wkhtmltopdf)
        file_name = f'{random.randint(10000, 99999)}_file.pdf'
        pdfkit.from_string(html_text, file_name, configuration=config)
        return file_name

    def get_all_assignment(self, request):
        try:
            data = AssignmentModel.objects.filter(
                user=request.user.id).order_by("-created_at")
            pagination_obj = CustomPagination()
            search_keys = []
            result = pagination_obj.custom_pagination(
                request, search_keys, categorySerializer.CreateAssignmentSerializers, data)
            return {"data": result, "message": messages.FETCH, "status": 200}
        except Exception as e:
            return {"data": None, "message": messages.WENT_WRONG, "status": 400}

    def update_download_file(self, request, id):
        try:
            assignment = AssignmentModel.objects.get(id=id)
            if request.data["type"] == 1:
                if request.data["new"] is True:
                    file = self.html_to_pdf(request)
                    return {"data": file, "message": "pdf generated successfully", "status": 200}
                if assignment.download_file:
                    return {"data": assignment.download_file, "message": messages.UPDATED, "status": 200}
                file = self.html_to_pdf(request)
                assignment.download_file = file
                assignment.save()
            if request.data["type"] == 2:
                if request.data["new"] is True:
                    file = self.html_to_doc(request)
                    return {"data": file, "message": "pdf generated successfully", "status": 200}
                if assignment.download_doc_file:
                    return {"data": assignment.download_doc_file, "message": messages.UPDATED, "status": 200}
                file = self.html_to_doc(request)
                assignment.download_doc_file = file
                assignment.save()
            elif request.data["type"] == 3:
                file = self.html_to_pdf(request)
                Thread(target=send_pdf_file_to_mail, args=(
                    assignment.user.email, file)).start()
                return {"data": None, "message": "File send to your email successfully", "status": 200}
            return {"data": file, "message": messages.UPDATED, "status": 200}
        except AssignmentModel.DoesNotExist:
            return {"data": None, "message": "Record not found", "status": 400}
        except Exception as e:
            return {"data": str(e), "message": messages.WENT_WRONG, "status": 400}

    def file_summary_download(self, request, id):
        try:
            file_summary = FileSumarizationModel.objects.get(id=id)
            if request.data["type"] == 2:
                if request.data.get("new"):
                    file = self.html_to_pdf(request)
                    file_summary.download_file = file
                    file_summary.save()
                elif not request.data.get("new") and file_summary.download_file:
                    return {"data": file_summary.download_file, "message": messages.UPDATED, "status": 200}
                else:
                    file = self.html_to_pdf(request)
                    file_summary.download_file = file
                    file_summary.save()
            elif request.data["type"] == 3:
                if file_summary.download_highlighted_file:
                    return {"data": file_summary.download_highlighted_file, "message": messages.UPDATED, "status": 200}
                else:
                    file = self.html_to_pdf(request)
                    file_summary.download_highlighted_file = file
                    file_summary.save()
            return {"data": file, "message": messages.UPDATED, "status": 200}
        except Exception as err:
            print(err, '========file error=========')
            return {"data": str(err), "message": messages.WENT_WRONG, "status": 400}

    def html_to_pdf(self, request):
        try:
            html_text = request.data["html_text"]
            # Update this path as necessary
            path_to_wkhtmltopdf = 'C:\\Program Files\\wkhtmltopdf\\bin\\wkhtmltopdf.exe'
            config = pdfkit.configuration(wkhtmltopdf=path_to_wkhtmltopdf)
            file_name = f'{random.randint(10000, 99999)}_file.pdf'
            pdfkit.from_string(html_text, file_name, configuration=config)
            saved_file = saveFile(file_name, "application/pdf")
            if os.path.exists(file_name):
                os.remove(file_name)
            return saved_file[0]
        except Exception as e:
            return {"data": None, "message": str(e), "status": status.HTTP_400_BAD_REQUEST}

    def html_to_doc(self, request):
        try:
            delete_files = []
            html_text = request.data["html_text"]
            # Update this path as necessary
            path_to_wkhtmltopdf = 'C:\\Program Files\\wkhtmltopdf\\bin\\wkhtmltopdf.exe'
            config = pdfkit.configuration(wkhtmltopdf=path_to_wkhtmltopdf)
            file_name = f'{random.randint(10000, 99999)}_file.pdf'
            doc_file_name = f'{random.randint(10000, 99999)}_file.docx'
            delete_files = [file_name, doc_file_name]
            pdfkit.from_string(html_text, file_name, configuration=config)
            cv = Converter(file_name)
            cv.convert(doc_file_name, start=0, end=None)
            cv.close()
            saved_file = saveFile(
                doc_file_name, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            for i in delete_files:
                if os.path.exists(i):
                    os.remove(i)
            return saved_file[0]
        except Exception as e:
            return {"data": None, "messages": str(e), "status": status.HTTP_400_BAD_REQUEST}

    def get_assignment_by_id(self, request, id):
        try:
            data = AssignmentModel.objects.get(id=id)
        except:
            return {"data": None, "message": messages.RECORD_NOT_FOUND, "status": 400}

        serializer = categorySerializer.CreateAssignmentSerializers(data)
        return {"data": serializer.data, "message": messages.FETCH, "status": 200}


############## settings app ###########


    def get_list_faq(self, request):
        try:
            user = FaqModel.objects.all()
        except:
            return {"data": None, "message": messages.RECORD_NOT_FOUND, "status": 400}
        serializer = adminSerializer.FaqModelSerializer(user, many=True)
        return {"data": serializer.data, "message": messages.FETCH, "status": 200}

    def get_terms_condition(self, request):
        try:
            user = CmsModel.objects.all()
        except:
            return {"data": None, "message": messages.RECORD_NOT_FOUND, "status": 400}
        serializer = adminSerializer.GetAllTermsConditionSerializer(
            user, many=True)
        return {"data": serializer.data, "message": messages.FETCH, "status": 200}

    def delete_user(self, request):
        try:
            user_obj = UserModel.objects.get(id=request.user.id)
        except UserModel.DoesNotExist:
            return {"data": None, "message": messages.RECORD_NOT_FOUND, "status": 400}
        user_obj.delete()
        return {"data": None, "message": messages.USER_DELETED, "status": 200}


# articles


    def get_article_response_list(self, request):

        topic = request.data.get("topic")
        words = request.data.get("words")
        language = request.data.get("language")
        region = request.data.get("region")
        tone = request.data.get("tone")
        pov = request.data.get("pronouns")

        llm = ChatGoogleGenerativeAI(model="gemini-pro")
        try:
            message_content_1 = f"Generate artcle topics list based on {topic} in {tone} tone of voice from a {pov} point of view in {language} language for a person from {region}. Output should contain only three topics headings(numbered like 1,2,3) and strictly two side headings(numbered like i, ii, iii)."
            message_content = message_content_1
            message = HumanMessage(
                content=[
                    {"type": "text", "text": message_content}
                ]
            )
            response = llm.invoke([message])
            result = to_markdown(response.content)

            save_article = ArticleModel.objects.create(
                user_id=request.user.id,
                topic=topic,
                language=language,
                region=region,
                pov=pov,
                words=words,
                tone=tone,
                result=result
            )

            return {"data": result, "record_id": save_article.id, "message": "Article generated successfully.", "status": 200}
        except Exception as e:
            return {"data": str(e), "message": messages.WENT_WRONG, "status": 400}

    def regenerate_article(self, request, id):
        request.data["record_id"] = id
        result = self.generate_detailed_article_based_on_topics(request)
        return result

    def generate_detailed_article_based_on_topics(self, request):
        try:
            llm = ChatGoogleGenerativeAI(model="gemini-pro")
            if "record_id" not in request.data:
                topic = request.data.get("topic")
                words = request.data.get("words")
                language = request.data.get("language").lower()
                region = request.data.get("region")
                tone = request.data.get("tone")
                pov = request.data.get("pronouns")
            elif "record_id" in request.data:
                get_article_record = ArticleModel.objects.get(
                    id=request.data.get("record_id"))
                topic = get_article_record.topic
                tone = get_article_record.tone
                pov = get_article_record.pov
                language = get_article_record.language.lower()
                region = get_article_record.region
                words = get_article_record.words
            ####
            if language == "english":
                QUERY = f"You are article generator. Generate an article on {topic} in the point of view of {pov} which I provide you. Whole Article should be in {language} and of approximately {words} words with voice of tone as {tone} and article should belongs to {region} region. Format should be descriptive. Strictly keep Headings(numbered as 1,2,3)."
                message_content = [
                    {
                        "type": "text",
                        "text": QUERY
                    }
                ]
                message = HumanMessage(content=message_content)
                response = llm.invoke([message])
                final_response = response.content.replace(
                    "*", "").replace("#", "").replace("-", "")
            elif language == "arabic":
                result = generate_article_util(topic, tone, pov, region, words)
                final_response = ""
                print(result, '----result-----result----')
                try:
                    for i, j in result["content"].items():
                        final_response += i + ". " + \
                            j["heading"] + "\n\n" + j["content"] + "\n\n"
                except:
                    return {"data": None, "message": "Please try again", "status": 400}
            # save record
            if "record_id" not in request.data:
                get_article_record = ArticleModel.objects.create(
                    user_id=request.user.id,
                    topic=topic,
                    language=language,
                    region=region,
                    pov=pov,
                    words=words,
                    tone=tone,
                    result=final_response
                )
            elif "record_id" in request.data:
                get_article_record.result = final_response
                get_article_record.save()
            print(final_response, '-----final_response----final_response-----')
            return {"data": final_response, "record_id": get_article_record.id, "message": "Detailed article generated successfully", "status": 200}
        except Exception as err:
            return {"data": str(err), "message": "Please try again", "status": 400}

    def get_article_history(self, request):
        try:
            articles = ArticleModel.objects.filter(
                user_id=request.user.id).order_by("-updated_at")
            pagination_obj = CustomPagination()
            search_keys = []
            result = pagination_obj.custom_pagination(
                request, search_keys, categorySerializer.GetArticlesListSerializer, articles)
            return {"data": result, "message": messages.FETCH, "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.WENT_WRONG, "status": 400}

    def get_article_by_id(self, request, id):
        try:
            article = ArticleModel.objects.get(id=id)
            serializer = categorySerializer.GetArticlesListSerializer(article)
            return {"data": serializer.data, "message": messages.FETCH, "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.WENT_WRONG, "status": 400}

    def download_article(self, request):
        try:
            if request.data.get("type") == 2:
                file = self.html_to_doc(request)
                return {"data": file, "message": "Document downloaded successfully", "status": 200}
            elif request.data.get("type") == 3:
                file = self.html_to_pdf(request)
                Thread(target=send_pdf_file_to_mail, args=(
                    request.user.email, file)).start()
                return {"data": None, "message": "Pdf sent to mail successfully", "status": 200}
            else:
                file = self.html_to_pdf(request)
                if isinstance(file, dict):
                    return {"data": file["message"], "message": messages.WENT_WRONG, "status": 400}
            return {"data": file, "message": messages.UPDATED, "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.WENT_WRONG, "status": 400}


### common for all ####


    def send_file_to_mail(self, request):
        try:
            data = UserModel.objects.get(id=request.user.id)
            email = data.email
            file_link = request.data["file_link"]
            sendMail.send_pdf_file_to_mail(email, file_link)

            return {"data": None, "message": messages.FILE_LINK_SEND, "status": 200}
        except Exception as e:
            return {"data": None, "message": messages.WENT_WRONG, "status": 400}

    def emu_to_pixels(self, emu):
        # Convert EMUs to pixels
        return int(emu * 96 / 914400)

    def draw_text(self, draw, text_frame, left, top):
        from PIL import Image, ImageDraw, ImageFont
        font = ImageFont.load_default()
        current_top = top
        for paragraph in text_frame.paragraphs:
            text = paragraph.text
            draw.text((left, current_top), text, fill='black', font=font)
            current_top += 15  # Move down for the next line

    def new_service(self, request):
        # import comtypes.client
        # try:
        #     powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        #     powerpoint.Visible = 1
        #     outputFileName = f"{random.randint(1000, 9999)}_outputFileName.pdf"
        #     deck = powerpoint.Presentations.Open("SykeIndia.pptx")
        #     deck.SaveAs(outputFileName, 32) # formatType = 32 for ppt to pdf
        #     deck.Close()
        #     powerpoint.Quit()

        #     return {"data": "", "message": "done", "status": 200}
        # except Exception as err:
        #     return {"data": str(err), "message": "went wrong", "status": 400}
        try:
            import aspose.slides as slides
            file = request.FILES.get("ppt_file")
            fs = FileSystemStorage()
            input_pdf_file = f"{random.randint(1000, 9999)}_{file.name}.pdf"
            fs.save(input_pdf_file, file)
            presentation = slides.Presentation(input_pdf_file)
            OUTPUT_FILE_NAME = generate_file_name(file.name)[0] + ".pdf"
            # Saves the presentation as a PDF
            output = f"{random.randint(1000, 9999)}_PPT-to-PDF.pdf"
            presentation.save(output, slides.export.SaveFormat.PDF)
            SAVED_FILE_RESPONSE = save_file_conversion(
                output, OUTPUT_FILE_NAME, "application/pdf")
            data = {
                "media_url": SAVED_FILE_RESPONSE[0],
                "media_type": "pdf",
                "media_name": SAVED_FILE_RESPONSE[1]
            }
            serializer = CreateUpdateUploadMediaSerializer(data=data)
            if serializer.is_valid():
                serializer.save()
            save_file_in_model = FileConversationModel.objects.create(
                user_id=request.user.id,
                converted_media_id=serializer.data["id"],
                sub_category=17
            )

            # Clean up temporary files
            if os.path.exists(output):
                os.remove(output)
            if os.path.exists(input_pdf_file):
                os.remove(input_pdf_file)

            return {
                "data": serializer.data,
                "message": messages.CONVERT_SUCCESS,
                "status": status.HTTP_200_OK
            }
        except Exception as err:
            return {"data": str(err), "message": "Something went wrong", "status": 400}

    def new_doc_to_pdf_service(self, request):
        from docx2pdf import convert

        try:
            file = request.FILES.get("word_file")
            generate_name = generate_file_name(file.name)
            FILE_NAME = generate_name[0]
            OUTPUT_FILE_NAME = generate_name[0] + ".pdf"
            name = f"{random.randint(1000, 9999)}_{FILE_NAME}"
            output_path = os.path.join(
                os.getcwd(), f"{random.randint(1000, 9999)}{name}.pdf")
            input_path = os.path.join(
                os.getcwd(), f"{random.randint(1000, 9999)}{name}.docx")
            fs = FileSystemStorage()
            fs.save(input_path, file)
            convert(input_path=input_path, output_path=output_path)
            SAVED_FILE_RESPONSE = save_file_conversion(
                output_path, OUTPUT_FILE_NAME, "application/pdf")
            data = {
                "media_url": SAVED_FILE_RESPONSE[0],
                "media_type": "pdf",
                "media_name": SAVED_FILE_RESPONSE[1]
            }
            serializer = CreateUpdateUploadMediaSerializer(data=data)
            if serializer.is_valid():
                serializer.save()
            save_file_in_model = FileConversationModel.objects.create(
                user_id=request.user.id,
                converted_media_id=serializer.data["id"],
                sub_category=11
            )
            if os.path.exists(input_path):
                os.remove(input_path)
            if os.path.exists(output_path):
                os.remove(output_path)
            return {"data": serializer.data, "message": "File converted successfully", "status": 200}
        except Exception as err:
            return {"data": str(err), "message": "Please upload the file again.", "status": 400}

    def ability(self, request):
        print(request.GET.get("type"),
              '---------------request.GET.get("type")-0------')
        try:
            api_type = True if request.GET.get("type") == "1" else False
            count = AbilityModel.objects.filter(
                is_arabic=False, is_mcq=api_type).count()
            if count <= 20:
                questions = AbilityModel.objects.filter(
                    is_arabic=False, is_mcq=api_type).order_by('?')
            else:
                random_ids = random.sample(list(AbilityModel.objects.filter(
                    is_arabic=False, is_mcq=api_type).values_list('id', flat=True)), 20)
                questions = AbilityModel.objects.filter(
                    id__in=random_ids, is_arabic=False, is_mcq=api_type)
            serialized_questions = [
                {
                    "question_no": idx + 1,
                    "question": question.question,
                    "answer_option": question.answer_option,
                    "correct_answer": question.corect_answer,
                    "answer": [question.corect_answer],
                    "user_answer": ""
                }
                for idx, question in enumerate(questions)
            ]
            if not serialized_questions:
                print(api_type is True, api_type is False)
                if api_type is True:
                    return {"data": messages.NO_QUESTIONS, "message": messages.NO_QUESTIONS, "status": 400}
                elif api_type is False:
                    return {"data": messages.NO_FLASHCARDS, "message": messages.NO_FLASHCARDS, "status": 400}
            save_record = TestingModel.objects.create(user_id=request.user.id,
                                                      sub_category=3,
                                                      sub_category_type=int(
                                                          request.GET.get("type")),
                                                      result=serialized_questions,
                                                      remaining_answers=len(serialized_questions))
            return {"data": serialized_questions, "record_id": save_record.id, "message": messages.FETCH, "status": 200}
        except Exception as e:
            return {"data": None, "message": messages.WENT_WRONG, "status": 400}

    def achievement(self, request, id):
        try:
            api_type = True if request.GET.get("type") == "1" else False
            print(api_type, '-----')
            count = AchievementModel.objects.filter(
                subject_id=id, is_arabic=False, is_mcq=api_type).count()
            if count <= 20:
                questions = AchievementModel.objects.filter(
                    subject_id=id, is_arabic=False, is_mcq=api_type).order_by('?')
            else:
                random_ids = AchievementModel.objects.filter(
                    subject_id=id, is_arabic=False, is_mcq=api_type).order_by('?').values_list('id', flat=True)[:20]
                questions = AchievementModel.objects.filter(
                    id__in=random_ids, is_arabic=False, is_mcq=api_type)
            serialized_questions = [
                {
                    "question_no": idx + 1,
                    "question": question.question,
                    "answer_option": question.answer_option,
                    "correct_answer": question.corect_answer,
                    "answer": [question.corect_answer],
                    "user_answer": ""
                }
                for idx, question in enumerate(questions)
            ]
            if not serialized_questions:
                if api_type is True:
                    return {"data": messages.NO_QUESTIONS, "message": messages.NO_QUESTIONS, "status": 400}
                elif api_type is False:
                    return {"data": messages.NO_FLASHCARDS, "message": messages.NO_FLASHCARDS, "status": 400}
            save_record = TestingModel.objects.create(user_id=request.user.id,
                                                      sub_category=4,
                                                      result=serialized_questions,
                                                      sub_category_type=int(
                                                          request.GET.get("type")),
                                                      remaining_answers=len(serialized_questions))
            return {"data": serialized_questions, "record_id": save_record.id, "message": messages.FETCH, "status": 200}
        except Exception as e:
            return {"data": str(e), "message": messages.WENT_WRONG, "status": 400}

    def get_testing_record_by_id(self, request, id):
        try:
            record = TestingModel.objects.get(id=id)
            serializer = categorySerializer.GetPreviousTestSerializer(record)
            return {"data": serializer.data, "message": messages.WENT_WRONG, "status": 200}
        except TestingModel.DoesNotExist:
            return {"data": None, "message": messages.RECORD_NOT_FOUND, "status": 400}
        except:
            return {"data": None, "message": messages.WENT_WRONG, "status": 400}

    def get_presentation_text(self, request):
        try:
            topic = request.data.get("topic")
            slides = request.data.get("slides")
            if "ar" in detect(topic):
                input_language = "arabic"
            else:
                input_language = "english"
            chatgpt_result = generate_presentation_util(
                topic, slides, input_language)
        # data = f"You are a presentation maker. Give me contents to make a presentation of {slides} slides on the topic - {topic} in {input_language} language. The content of each slide should be more than 15000 words strictly with proper headings. So fill up the content part in pointers. Format should be in python json dictionary and keys should be strictly (number, heading ,content(the matter in content should be around 150-200 words for each slide strictly))"
        # # data=f"You are a topics list generator. Generate research topics list based on {topic}. Output should contain only three topics headings(numbered like 1,2,3) and strictly two side headings(numbered like i, ii, iii)."
        # query = data
        # llm = ChatGoogleGenerativeAI(model="gemini-pro")
        # try:
        #     response = llm.invoke(query)
        #     result = to_markdown(response.content)
        #     print(result, '-----')
        #     result = result.replace("Slide ", '')
        #     reversed_result = result[::-1]
        #     response = result[result.index(
        #         "{"):len(result) - (reversed_result.index("}")) + 1]
        #     final_response = ast.literal_eval(response)
        #     final_response = list(map(lambda v: v, final_response.values()))
            chatgpt_final_response = list(
                map(lambda v: v, chatgpt_result.values()))
            return {"data": chatgpt_final_response, "message": messages.PRESENTATION_GENERATED, "status": 200}
        except Exception as err:
            print(err, type(err), '-------------err-----------')
            return {"data": None, "message": messages.TRY_AGAIN, "status": 400}

    def save_presentation_binary_data(self, request):
        try:
            presentation_obj = PresentationModel()
            presentation_obj.user_id = request.user.id
            presentation_obj.template_id = request.data.get("template_id", 1)
            presentation_obj.slides = request.data.get("slides", 6)
            presentation_obj.text = request.data.get("topic")
            presentation_obj.binary_data = request.data.get("binary_data")
            presentation_obj.save()
            return {"data": {"record_id": presentation_obj.id}, "message": "Binary data saved successfully", "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.TRY_AGAIN, "status": 400}

    def update_presentation_by_id(self, request, id):
        try:
            presentation_obj = PresentationModel.objects.get(id=id)
            presentation_obj.binary_data = request.data.get("binary_data")
            presentation_obj.template_id = request.data.get("template_id")
            presentation_obj.save()
            return {"data": [], "message": "Presentation updated successfully", "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.TRY_AGAIN, "status": 400}

    def get_presentation_by_id(self, request, id):
        try:
            presentation_obj = PresentationModel.objects.get(id=id)
            data = {}
            data["id"] = presentation_obj.id
            data["slides"] = presentation_obj.slides
            data["text"] = presentation_obj.text
            data["template_id"] = presentation_obj.template_id
            data["binary_data"] = ast.literal_eval(
                presentation_obj.binary_data)
            return {"data": data, "message": "Presentation fetched successfully", "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.TRY_AGAIN, "status": 400}

    def presentation_history(self, request):
        try:
            presentation_objs = PresentationModel.objects.filter(
                user=request.user).order_by("-updated_at")
            pagination_obj = CustomPagination()
            search_keys = []
            result = pagination_obj.custom_pagination(
                request, search_keys, categorySerializer.PresentationHistorySerializer, presentation_objs)
            return {"data": result, "message": "Presentation history fetched successfully", "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.TRY_AGAIN, "status": 400}

    def save_notes(self, request):
        save_notes = NoteTakingModel.objects.create(
            user_id=request.user.id, **request.data)
        return {"data": None, "message": messages.NOTES_ADDED, "status": 200}

    def notes_history(self, request):
        try:
            if "filter" not in request.data:
                all_notes = NoteTakingModel.objects.filter(
                    user=request.user).order_by("-created_at")
            elif "filter" in request.data:
                all_notes = NoteTakingModel.objects.filter(
                    user=request.user, created_at__date=request.data["filter"]).order_by("-created_at")
            pagination_obj = CustomPagination()
            search_keys = []
            result = pagination_obj.custom_pagination(
                request, search_keys, categorySerializer.AllNotesSerializer, all_notes)
            try:
                today_date = datetime.now()
                recent_notes = []
                if "filter" not in request.data:
                    for i in range(10):
                        date_wise_notes = NoteTakingModel.objects.filter(created_at__date=today_date.date(), user=request.user).values(
                            "id", "type", "note_screenshot", "canvas_height")
                        notes_dict = {"title": datetime.strftime(
                            today_date, "%Y-%m-%d"), "count": len(date_wise_notes)}
                        today_date -= timedelta(days=1)
                        recent_notes.append(notes_dict)
            except Exception as err:
                print(err)
            return {"data": result, "recent_notes": recent_notes, "message": messages.NOTES_HISTORY, "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.TRY_AGAIN, "status": 400}

    def notes_by_id(self, request, id):
        try:
            notes = NoteTakingModel.objects.get(id=id)
            serializer = categorySerializer.NoteTakingSerializer(notes)
            return {"data": serializer.data, "message": messages.NOTES_FETCHED, "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.WENT_WRONG, "status": 400}

    def edit_notes_by_id(self, request, id):
        try:
            notes = NoteTakingModel.objects.get(id=id)
            if "canvas_height" in request.data:
                notes.canvas_height = request.data["canvas_height"]
            if "comments" in request.data:
                notes.comments = request.data["comments"]
            if "text_timestamp" in request.data:
                notes.text_timestamp = request.data["text_timestamp"]
            if "binary_data" in request.data:
                notes.binary_data = request.data["binary_data"]
            if "note_screenshot" in request.data:
                notes.note_screenshot = request.data["note_screenshot"]
            notes.save()
            return {"data": {}, "message": messages.NOTES_UPDATED, "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.WENT_WRONG, "status": 400}

    def notes_actions(self, request):
        try:
            if request.data.get("type") == 1:
                message = messages.NOTES_DUPLICATED
                records = NoteTakingModel.objects.filter(
                    id__in=request.data.get("record_ids"), user=request.user)
                for i in records:
                    save_notes = NoteTakingModel.objects.create(
                        user_id=request.user.id,
                        type=i.type,
                        binary_data=i.binary_data,
                        note_screenshot=i.note_screenshot,
                        canvas_height=i.canvas_height,
                        is_duplicate=True
                    )
            elif request.data.get("type") == 2:
                message = messages.NOTES_FAVOURITED
                records = NoteTakingModel.objects.filter(
                    id__in=request.data.get("record_ids"), user=request.user)
                for i in records:
                    if i.is_favourite is True:
                        i.is_favourite = False
                    else:
                        i.is_favourite = True
                    i.save()
            elif request.data.get("type") == 3:
                message = messages.NOTES_DELETED
                records = NoteTakingModel.objects.filter(
                    id__in=request.data.get("record_ids"), user=request.user)
                records.delete()
            return {"data": None, "message": message, "status": 200}
        except Exception as err:
            return {"data": str(err), "message": messages.WENT_WRONG, "status": 400}
